"""
Gera apresentação executiva com python-pptx a partir de MAIO 2026 - Copia.xlsx (read-only).
Paleta corporativa: fundo claro, textos grafite/azul-marinho, verde/laranja para destaques.
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from danone import (
    EstudoDanone,
    LinhaFaturamento,
    PLANILHA_FONTE,
    carregar_estudo,
    fmt_moeda,
    fmt_ms,
    fmt_pct,
    fmt_produto_curto,
)

# Alias legado
DadosApresentacao = EstudoDanone
carregar_dados = carregar_estudo

# --- Paleta corporativa ---
COLOR_BG = RGBColor(248, 249, 250)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_NAVY = RGBColor(26, 43, 74)
COLOR_GRAPHITE = RGBColor(44, 62, 80)
COLOR_MUTED = RGBColor(108, 117, 125)
COLOR_GREEN_BG = RGBColor(232, 245, 233)
COLOR_GREEN_TEXT = RGBColor(27, 94, 32)
COLOR_ORANGE_BG = RGBColor(255, 243, 224)
COLOR_ORANGE_TEXT = RGBColor(191, 87, 0)
COLOR_BAR_2025 = "#6C7A89"
COLOR_BAR_2026 = "#1A2B4A"
COLOR_TEAL = "#2E86AB"

FONT_BODY = "Segoe UI"
COLOR_LINK = RGBColor(0, 86, 179)

from danone.config import PASTA_APRESENTACOES

PASTA_ENTREGA = PASTA_APRESENTACOES
PPTX_LOCAL = "Apresentacao MAIO 2026.pptx"
PPTX_EMAIL = "Apresentacao MAIO 2026_email.pptx"
EXCEL_LINK_LABEL = "Abrir planilha de dados (Excel)"


def _set_slide_background(slide, color: RGBColor = COLOR_BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_GRAPHITE,
    align=PP_ALIGN.LEFT,
    font_name: str = FONT_BODY,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color
    p.alignment = align
    return box


def _add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    font_size: int = 16,
    color: RGBColor = COLOR_GRAPHITE,
    spacing_after: int = 10,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.name = FONT_BODY
        p.font.color.rgb = color
        p.space_after = Pt(spacing_after)
    return box


def _hyperlink_excel(excel_path: Path, modo: str = "local") -> str:
    path = excel_path.resolve()
    if modo == "email":
        return path.name
    return path.as_uri()


def _caminhos_entrega(modo_link: str, planilha_fonte: Path) -> tuple[Path, Path]:
    PASTA_ENTREGA.mkdir(parents=True, exist_ok=True)
    pptx = PASTA_ENTREGA / (PPTX_EMAIL if modo_link == "email" else PPTX_LOCAL)
    return pptx, planilha_fonte.resolve()


def _add_hyperlink_run(paragraph, text: str, address: str, font_size: int = 11) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = FONT_BODY
    run.font.bold = True
    run.font.underline = True
    run.font.color.rgb = COLOR_LINK
    run.hyperlink.address = address


def _add_rodape_base_dados(slide, excel_address: str, periodo: str) -> None:
    rodape_top = Inches(6.95)
    line = slide.shapes.add_shape(
        1, Inches(0.6), rodape_top, Inches(12.1), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(222, 226, 230)
    line.line.fill.background()

    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.38)
    )
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    run_prefix = p.add_run()
    run_prefix.text = f"{periodo}  ·  Base analítica  ·  "
    run_prefix.font.size = Pt(11)
    run_prefix.font.name = FONT_BODY
    run_prefix.font.color.rgb = COLOR_MUTED

    _add_hyperlink_run(p, EXCEL_LINK_LABEL, excel_address, font_size=11)


def _add_callout_excel_capa(slide, excel_address: str) -> None:
    left, top = Inches(0.8), Inches(5.05)
    width, height = Inches(7.8), Inches(0.95)

    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = RGBColor(222, 226, 230)
    card.line.width = Pt(0.75)

    accent = slide.shapes.add_shape(1, left, top, Inches(0.08), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_NAVY
    accent.line.fill.background()

    box = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.12), width - Inches(0.4), height - Inches(0.2)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    run_t = p.add_run()
    run_t.text = "Aprofundamento executivo  ·  "
    run_t.font.size = Pt(14)
    run_t.font.name = FONT_BODY
    run_t.font.bold = True
    run_t.font.color.rgb = COLOR_NAVY

    _add_hyperlink_run(p, EXCEL_LINK_LABEL, excel_address, font_size=14)

    p2 = tf.add_paragraph()
    run_d = p2.add_run()
    run_d.text = (
        "Planilha ESTUDO_DANONE_MAT_MAIO com ranking, regiões, bandeiras ABRAFAD e base detalhada. "
        "Por e-mail: salve o .pptx e o .xlsx na mesma pasta antes de clicar."
    )
    run_d.font.size = Pt(12)
    run_d.font.name = FONT_BODY
    run_d.font.bold = False
    run_d.font.color.rgb = COLOR_MUTED


def _criar_grafico_barras_labs(labs: list[LinhaFaturamento]) -> io.BytesIO:
    nomes = [lab.nome.replace(" BABY NUTRIT", "\nBABY NUTRIT") for lab in labs]
    valores_2025 = [lab.fat_2025 / 1e9 for lab in labs]
    valores_2026 = [lab.fat_2026 / 1e9 for lab in labs]
    crescimentos = [fmt_pct(lab.crescimento) for lab in labs]
    shares = [fmt_ms(lab.market_share) for lab in labs]

    x = range(len(nomes))
    width = 0.35

    fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=150)
    fig.patch.set_facecolor("#F8F9FA")
    ax.set_facecolor("#F8F9FA")

    bars_2025 = ax.bar(
        [i - width / 2 for i in x],
        valores_2025,
        width,
        label="Abr/25",
        color=COLOR_BAR_2025,
        edgecolor="none",
    )
    bars_2026 = ax.bar(
        [i + width / 2 for i in x],
        valores_2026,
        width,
        label="Abr/26",
        color=COLOR_BAR_2026,
        edgecolor="none",
    )

    ax.set_ylabel("Faturamento (R$ bilhões)", fontsize=11, color="#2C3E50")
    ax.set_xticks(list(x))
    ax.set_xticklabels(nomes, fontsize=10, color="#2C3E50")
    ax.tick_params(axis="y", labelsize=10, colors="#2C3E50")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#DEE2E6")
    ax.spines["bottom"].set_color("#DEE2E6")
    ax.grid(axis="y", linestyle="--", alpha=0.35, color="#ADB5BD")
    ax.legend(frameon=False, fontsize=10, loc="upper right")

    for bar, pct, ms in zip(bars_2026, crescimentos, shares):
        h = bar.get_height()
        cor = "#1B5E20" if not pct.startswith("-") else "#BF5700"
        ax.annotate(
            pct,
            xy=(bar.get_x() + bar.get_width() / 2, h),
            xytext=(0, 5),
            textcoords="offset points",
            ha="center",
            va="bottom",
            fontsize=9,
            fontweight="bold",
            color=cor,
        )
        if ms != "—":
            ax.annotate(
                f"MS {ms}",
                xy=(bar.get_x() + bar.get_width() / 2, h),
                xytext=(0, -14),
                textcoords="offset points",
                ha="center",
                va="top",
                fontsize=8,
                color="#6C757D",
            )

    plt.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    buffer.seek(0)
    return buffer


def _criar_grafico_bandeiras_brasil(bandeiras: list[LinhaFaturamento]) -> io.BytesIO:
    top = bandeiras[:8]
    nomes = [b.nome.title()[:28] for b in top]
    valores = [b.fat_2026 / 1e6 for b in top]
    shares = [fmt_ms(b.market_share) for b in top]
    cresc = [fmt_pct(b.crescimento) for b in top]

    fig, ax = plt.subplots(figsize=(7.0, 4.0), dpi=150)
    fig.patch.set_facecolor("#F8F9FA")
    ax.set_facecolor("#F8F9FA")

    bars = ax.barh(nomes, valores, color=COLOR_TEAL, edgecolor="none", height=0.55)
    ax.set_xlabel("Faturamento Abr/26 (R$ milhões)", fontsize=11, color="#2C3E50")
    ax.tick_params(axis="both", labelsize=9, colors="#2C3E50")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#DEE2E6")
    ax.spines["bottom"].set_color("#DEE2E6")
    ax.grid(axis="x", linestyle="--", alpha=0.35, color="#ADB5BD")

    for bar, val, ms, yoy in zip(bars, valores, shares, cresc):
        ax.text(
            val + max(valores) * 0.02,
            bar.get_y() + bar.get_height() / 2,
            f"{yoy} · MS {ms}",
            va="center",
            ha="left",
            fontsize=8,
            fontweight="bold",
            color="#2C3E50",
        )

    plt.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    buffer.seek(0)
    return buffer


def _criar_grafico_regioes(regioes: list[LinhaFaturamento]) -> io.BytesIO:
    nomes = [r.nome.title() for r in regioes]
    cresc = [(r.crescimento or 0) * 100 for r in regioes]
    cores = ["#1B5E20" if v >= 0 else "#BF5700" for v in cresc]

    fig, ax = plt.subplots(figsize=(7.0, 3.8), dpi=150)
    fig.patch.set_facecolor("#F8F9FA")
    ax.set_facecolor("#F8F9FA")

    bars = ax.barh(nomes, cresc, color=cores, edgecolor="none", height=0.55)
    ax.axvline(0, color="#ADB5BD", linewidth=0.8)
    ax.set_xlabel("Crescimento YoY (%)", fontsize=11, color="#2C3E50")
    ax.tick_params(axis="both", labelsize=10, colors="#2C3E50")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#DEE2E6")
    ax.spines["bottom"].set_color("#DEE2E6")
    ax.grid(axis="x", linestyle="--", alpha=0.35, color="#ADB5BD")

    for bar, val, reg in zip(bars, cresc, regioes):
        offset = 0.4 if val >= 0 else -0.4
        ha = "left" if val >= 0 else "right"
        ms_txt = fmt_ms(reg.market_share)
        label = f"{val:+.1f}%".replace(".", ",")
        if ms_txt != "—":
            label += f" · MS {ms_txt}"
        ax.text(
            val + offset,
            bar.get_y() + bar.get_height() / 2,
            label,
            va="center",
            ha=ha,
            fontsize=8,
            fontweight="bold",
            color="#2C3E50",
        )

    plt.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    buffer.seek(0)
    return buffer


def _bullets_visao_geral(dados: DadosApresentacao) -> list[str]:
    bullets: list[str] = []
    if dados.laboratorios_destaque:
        top_cresc = max(dados.laboratorios_destaque, key=lambda x: x.crescimento or -999)
        bullets.append(
            f"{top_cresc.nome} lidera o crescimento entre os principais players "
            f"({fmt_pct(top_cresc.crescimento)})."
        )
        maior_fat = max(dados.laboratorios_destaque, key=lambda x: x.fat_2026)
        ms_txt = (
            f" e {fmt_ms(maior_fat.market_share)} de market share"
            if maior_fat.market_share is not None
            else ""
        )
        bullets.append(
            f"{maior_fat.nome} mantém a maior fatia de receita "
            f"({fmt_moeda(maior_fat.fat_2026)}{ms_txt})."
        )
    retrações = [l for l in dados.laboratorios_destaque if (l.crescimento or 0) < 0]
    if retrações:
        lab = retrações[0]
        bullets.append(
            f"{lab.nome} apresenta retração de {fmt_pct(lab.crescimento, com_sinal=False)} "
            f"no período, demandando revisão de estratégia."
        )
    if dados.total_danone and dados.total_danone.crescimento is not None:
        bullets.append(
            f"Portfólio Danone (NAO_MEDICAMENTO_NTR) cresce "
            f"{fmt_pct(dados.total_danone.crescimento)} consolidado."
        )
    return bullets[:4]


def _slide_capa(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(11.5),
        Inches(1.6),
        f"Desempenho de Categoria – Faturamento YoY\n({dados.periodo_label})",
        font_size=38,
        bold=True,
        color=COLOR_NAVY,
    )
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(3.85),
        Inches(11),
        Inches(0.8),
        "Análise de Impactos, Market Share, Laboratórios, Bandeiras Brasil e Regiões",
        font_size=22,
        color=COLOR_MUTED,
    )

    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.85), Inches(2.2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_NAVY
    line.line.fill.background()

    if dados.total_danone:
        kpi = (
            f"Danone NTR · Abr/25: {fmt_moeda(dados.total_danone.fat_2025)}  →  "
            f"Abr/26: {fmt_moeda(dados.total_danone.fat_2026)}  "
            f"({fmt_pct(dados.total_danone.crescimento)})"
        )
        _add_textbox(
            slide, Inches(0.8), Inches(4.55), Inches(11), Inches(0.45),
            kpi, font_size=14, bold=True, color=COLOR_GRAPHITE,
        )

    _add_callout_excel_capa(slide, excel_address)
    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def _slide_visao_geral(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
        "Visão Geral – Top 3 Laboratórios",
        font_size=36, bold=True, color=COLOR_NAVY,
    )

    labs = dados.laboratorios_destaque[:3]
    chart_buffer = _criar_grafico_barras_labs(labs)
    slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(1.15), width=Inches(7.2))

    _add_bullets(
        slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(4.5),
        _bullets_visao_geral(dados), font_size=15,
    )

    nota_linhas = [
        f"{lab.nome}: Abr/25 = {fmt_moeda(lab.fat_2025)} | Abr/26 = {fmt_moeda(lab.fat_2026)}"
        f" | MS = {fmt_ms(lab.market_share)} | YoY = {fmt_pct(lab.crescimento)}"
        for lab in labs
    ]
    _add_textbox(
        slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.55),
        "\n".join(nota_linhas), font_size=11, color=COLOR_MUTED,
    )
    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def _slide_regional(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
        "Performance Regional – Danone NTR",
        font_size=36, bold=True, color=COLOR_NAVY,
    )

    if dados.regioes:
        chart = _criar_grafico_regioes(dados.regioes)
        slide.shapes.add_picture(chart, Inches(0.55), Inches(1.2), width=Inches(6.8))

    bullets: list[str] = []
    if dados.regioes:
        lider = max(dados.regioes, key=lambda x: x.crescimento or -999)
        maior_fat = max(dados.regioes, key=lambda x: x.fat_2026)
        bullets.append(
            f"{lider.nome.title()} lidera crescimento regional ({fmt_pct(lider.crescimento)}"
            f"{', MS ' + fmt_ms(lider.market_share) if lider.market_share is not None else ''})."
        )
        bullets.append(
            f"{maior_fat.nome.title()} concentra maior faturamento "
            f"({fmt_moeda(maior_fat.fat_2026)} em Abr/26"
            f"{'; MS ' + fmt_ms(maior_fat.market_share) if maior_fat.market_share is not None else ''})."
        )
    if dados.total_regioes:
        bullets.append(
            f"Total Brasil: {fmt_moeda(dados.total_regioes.fat_2025)} → "
            f"{fmt_moeda(dados.total_regioes.fat_2026)} "
            f"({fmt_pct(dados.total_regioes.crescimento)})."
        )
    bullets.append("Priorizar expansão nas regiões de maior tração (Norte e Centro-Oeste).")

    _add_bullets(
        slide, Inches(7.6), Inches(1.5), Inches(5.2), Inches(4.8),
        bullets, font_size=15,
    )
    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def _slide_ranking_brasil(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
        "Ranking Brasil – Principais Bandeiras ABRAFAD",
        font_size=36, bold=True, color=COLOR_NAVY,
    )

    bandeiras = dados.bandeiras[:10]
    if bandeiras:
        chart = _criar_grafico_bandeiras_brasil(bandeiras)
        slide.shapes.add_picture(chart, Inches(0.55), Inches(1.15), width=Inches(6.9))

    bullets: list[str] = []
    if bandeiras:
        lider = bandeiras[0]
        bullets.append(
            f"{lider.nome.title()} lidera entre as bandeiras mapeadas "
            f"({fmt_moeda(lider.fat_2026)} em Abr/26; MS {fmt_ms(lider.market_share)})."
        )
        maior_ms = max(
            (b for b in bandeiras if b.market_share is not None),
            key=lambda x: x.market_share or 0,
            default=None,
        )
        if maior_ms:
            bullets.append(
                f"Maior participação individual: {maior_ms.nome.title()} "
                f"({fmt_ms(maior_ms.market_share)} de market share)."
            )
    if dados.total_bandeiras_brasil and dados.total_bandeiras_brasil.market_share:
        bullets.append(
            f"Demais concorrentes (agregado): MS {fmt_ms(dados.total_bandeiras_brasil.market_share)} "
            f"no universo de bandeiras."
        )
    elif dados.concorrentes:
        conc = next((c for c in dados.concorrentes if c.nome.upper() == "CONCORRENTES"), None)
        if conc:
            bullets.append(
                f"Canal concorrente agregado: {fmt_moeda(conc.fat_2026, True)} "
                f"({fmt_pct(conc.crescimento)} no período)."
            )
    bullets.append(
        "Priorizar parcerias nas redes de maior tração e share incremental no período."
    )

    _add_bullets(
        slide, Inches(7.5), Inches(1.4), Inches(5.3), Inches(4.5),
        bullets[:4], font_size=15,
    )

    if bandeiras:
        notas = [
            f"{b.nome.title()}: {fmt_moeda(b.fat_2026)} · MS {fmt_ms(b.market_share)} · {fmt_pct(b.crescimento)}"
            for b in bandeiras[:5]
        ]
        _add_textbox(
            slide, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.75),
            "\n".join(notas), font_size=10, color=COLOR_MUTED,
        )

    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def _slide_portfolio_danone(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
        "Portfólio Danone – Unidades de Negócio",
        font_size=36, bold=True, color=COLOR_NAVY,
    )

    linhas = ["Unidade | Abr/25 | Abr/26 | YoY | Market Share %"]
    for item in dados.portfolio_danone:
        linhas.append(
            f"{item.nome} | {fmt_moeda(item.fat_2025)} | {fmt_moeda(item.fat_2026)} | "
            f"{fmt_pct(item.crescimento)} | {fmt_ms(item.market_share)}"
        )
    if dados.total_danone:
        linhas.append(
            f"TOTAL | {fmt_moeda(dados.total_danone.fat_2025)} | "
            f"{fmt_moeda(dados.total_danone.fat_2026)} | "
            f"{fmt_pct(dados.total_danone.crescimento)} | {fmt_ms(dados.total_danone.market_share)}"
        )

    _add_textbox(
        slide, Inches(0.7), Inches(1.3), Inches(11.8), Inches(3.2),
        "\n".join(linhas), font_size=15, color=COLOR_GRAPHITE,
    )

    bullets = []
    baby = next((p for p in dados.portfolio_danone if "BABY" in p.nome.upper()), None)
    med = next((p for p in dados.portfolio_danone if "MEDICAL" in p.nome.upper()), None)
    if baby:
        bullets.append(
            f"Baby Nutrit é o principal motor (+{fmt_pct(baby.crescimento, com_sinal=False)}), "
            f"representando a maior fatia do portfólio."
        )
    if med and (med.crescimento or 0) < 0:
        bullets.append(
            f"Medical Nut apresenta leve retração ({fmt_pct(med.crescimento)}); "
            f"monitorar mix e canais."
        )
    if dados.total_danone:
        bullets.append(
            f"Crescimento consolidado de {fmt_pct(dados.total_danone.crescimento)} "
            f"supera a média de vários concorrentes diretos."
        )

    _add_bullets(
        slide, Inches(0.7), Inches(4.6), Inches(11.5), Inches(2.0),
        bullets, font_size=16,
    )
    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def _add_colored_block(
    slide, left, top, width, height,
    bg_color: RGBColor, title: str, produtos: str,
    insight_label: str, insight_text: str, title_color: RGBColor,
):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(222, 226, 230)
    shape.line.width = Pt(0.75)

    margin = Inches(0.25)
    inner_left = left + margin
    inner_width = width - 2 * margin

    _add_textbox(slide, inner_left, top + Inches(0.2), inner_width, Inches(0.45),
                 title, font_size=20, bold=True, color=title_color)
    _add_textbox(slide, inner_left, top + Inches(0.75), inner_width, Inches(1.1),
                 produtos, font_size=15, color=COLOR_GRAPHITE)
    _add_textbox(slide, inner_left, top + Inches(1.95), inner_width, Inches(0.35),
                 insight_label, font_size=16, bold=True, color=title_color)
    _add_textbox(slide, inner_left, top + Inches(2.35), inner_width, Inches(1.5),
                 insight_text, font_size=15, color=COLOR_GRAPHITE)


def _texto_produtos_impacto(produtos: list[LinhaFaturamento]) -> str:
    partes = []
    for p in produtos[:3]:
        nome = fmt_produto_curto(p.nome)
        pct = fmt_pct(p.crescimento)
        partes.append(f"{nome} ({pct})")
    return "Produtos: " + ", ".join(partes) + "."


def _insight_positivos(pos: list[LinhaFaturamento]) -> str:
    if not pos:
        return "Sem destaques positivos identificados no período."
    lider = pos[0]
    nome = fmt_produto_curto(lider.nome)
    delta = lider.fat_2026 - lider.fat_2025
    if (lider.crescimento or 0) > 1.0:
        return (
            f"{nome} registra crescimento excepcional de {fmt_pct(lider.crescimento)}, "
            f"com incremento de {fmt_moeda(delta)} no faturamento. "
            f"A linha Aptamil Profutura também mantém forte tração (+98% e +79%)."
        )
    return (
        "A linha Aptamil Profutura consolida-se como principal motor de tração, "
        "com forte incremento incremental de receita."
    )


def _insight_negativos(neg: list[LinhaFaturamento]) -> str:
    if not neg:
        return "Sem alertas relevantes no período."
    pior = neg[0]
    perda = pior.fat_2025 - pior.fat_2026
    return (
        f"{fmt_produto_curto(pior.nome)} lidera a retração ({fmt_pct(pior.crescimento)}), "
        f"com perda de {fmt_moeda(perda)}. Contração na Milnutri também demanda "
        f"ação imediata para estancar perda de receita."
    )


def _slide_extremos(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
        "Os Extremos – Impactos Positivos e Negativos",
        font_size=36, bold=True, color=COLOR_NAVY,
    )

    pos = dados.impactos_positivos
    neg = dados.impactos_negativos

    _add_colored_block(
        slide, Inches(0.55), Inches(1.2), Inches(5.9), Inches(4.2),
        COLOR_GREEN_BG, "Destaques Positivos",
        _texto_produtos_impacto(pos),
        "Insight:",
        _insight_positivos(pos),
        COLOR_GREEN_TEXT,
    )

    _add_colored_block(
        slide, Inches(6.75), Inches(1.2), Inches(5.9), Inches(4.2),
        COLOR_ORANGE_BG, "Pontos de Alerta",
        _texto_produtos_impacto(neg),
        "Insight:",
        _insight_negativos(neg),
        COLOR_ORANGE_TEXT,
    )
    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def _slide_plano_acao(prs: Presentation, excel_address: str, dados: DadosApresentacao) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)

    _add_textbox(
        slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
        "Recomendações Estratégicas",
        font_size=36, bold=True, color=COLOR_NAVY,
    )

    items: list[str] = []
    if pos := dados.impactos_positivos:
        rr = pos[0]
        items.append(
            f"Capitalizar o crescimento excepcional de {fmt_produto_curto(rr.nome)} "
            f"({fmt_pct(rr.crescimento)}) — garantir abastecimento e visibilidade em PDV."
        )
        prof = next((p for p in pos if "PROFUTURA" in p.nome.upper()), None)
        if prof:
            items.append(
                f"Reforçar trade marketing na linha Aptamil Profutura "
                f"({fmt_pct(prof.crescimento)} no principal SKU)."
            )
    else:
        items.append(
            "Garantir abastecimento e reforçar trade marketing na linha Aptamil Profutura."
        )
    if neg := dados.impactos_negativos:
        ar = next((p for p in neg if "AR" in p.nome.upper()), neg[0])
        items.append(
            f"Investigar imediatamente a queda de {fmt_pct(ar.crescimento, com_sinal=False)} "
            f"em {fmt_produto_curto(ar.nome)} (estoque, preço ou mix de canal)."
        )
        mil = next((p for p in neg if "MILNUTRI" in p.nome.upper()), None)
        if mil:
            items.append(
                f"Revisar posicionamento e promoção da linha Milnutri "
                f"({fmt_pct(mil.crescimento)} no período)."
            )
    if dados.regioes:
        lider = max(dados.regioes, key=lambda x: x.crescimento or -999)
        items.append(
            f"Intensificar ações comerciais em {lider.nome.title()} "
            f"(maior crescimento regional: {fmt_pct(lider.crescimento)})."
        )
    items.append(
        "Manter monitoramento via ESTUDO_DANONE_MAT_MAIO "
        "(Ranking, ABRAFAD, Região e Base detalhada)."
    )
    if dados.bandeiras:
        top_bandeira = dados.bandeiras[0]
        items.append(
            f"Fortalecer presença em {top_bandeira.nome.title()} "
            f"(líder entre bandeiras: MS {fmt_ms(top_bandeira.market_share)})."
        )

    _add_bullets(
        slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
        items[:5], font_size=17, spacing_after=14,
    )
    _add_rodape_base_dados(slide, excel_address, dados.periodo_label)


def gerar_apresentacao(
    caminho_saida: Path | None = None,
    caminho_excel: Path | None = None,
    planilha_fonte: Path | None = None,
    modo_link: str = "local",
) -> tuple[Path, Path]:
    """
    modo_link:
      - 'local'  -> file:/// (teste no seu PC; abre Excel direto)
      - 'email'  -> nome do arquivo (melhor para quem receber os dois anexos)
    """
    fonte = Path(planilha_fonte or PLANILHA_FONTE)
    if not fonte.exists():
        raise FileNotFoundError(f"Planilha base não encontrada: {fonte}")

    if caminho_saida is None and caminho_excel is None:
        caminho_pptx, caminho_xlsx = _caminhos_entrega(modo_link, fonte)
    else:
        caminho_pptx = caminho_saida or _caminhos_entrega(modo_link, fonte)[0]
        caminho_xlsx = caminho_excel or fonte.resolve()

    print(f"Lendo planilha (read-only): {fonte}")
    dados = carregar_estudo(fonte)

    print("\n--- Conferencia de dados (impactos) ---")
    for p in dados.impactos_positivos:
        print(f"  [+] {fmt_produto_curto(p.nome)}: {fmt_pct(p.crescimento)} | "
              f"{fmt_moeda(p.fat_2025)} -> {fmt_moeda(p.fat_2026)}")
    for p in dados.impactos_negativos:
        print(f"  [-] {fmt_produto_curto(p.nome)}: {fmt_pct(p.crescimento)} | "
              f"{fmt_moeda(p.fat_2025)} -> {fmt_moeda(p.fat_2026)}")
    print("---\n")

    excel_link = _hyperlink_excel(caminho_xlsx, modo=modo_link)
    print(f"Link da planilha ({modo_link}): {excel_link}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _slide_capa(prs, excel_link, dados)
    _slide_visao_geral(prs, excel_link, dados)
    _slide_regional(prs, excel_link, dados)
    _slide_ranking_brasil(prs, excel_link, dados)
    _slide_portfolio_danone(prs, excel_link, dados)
    _slide_extremos(prs, excel_link, dados)
    _slide_plano_acao(prs, excel_link, dados)

    prs.save(str(caminho_pptx))
    caminho_resolvido = caminho_pptx.resolve()
    print(f"\nApresentacao salva: {caminho_resolvido}")
    print(f"Planilha base (nao alterada): {caminho_xlsx}")

    # Marcador visivel no Cursor (sidebar as vezes nao mostra .pptx)
    marcador = PASTA_ENTREGA / "APRESENTACAO_GERADA.txt"
    marcador.write_text(
        f"Apresentacao gerada com sucesso.\n\n"
        f"Arquivo: {caminho_resolvido}\n"
        f"Planilha: {caminho_xlsx}\n"
        f"Slides: 7\n",
        encoding="utf-8",
    )
    print(f"Marcador: {marcador.resolve()}")

    return caminho_pptx, caminho_xlsx


if __name__ == "__main__":
    import sys

    modo = "email" if len(sys.argv) > 1 and sys.argv[1].lower() == "email" else "local"

    print(f"Pasta de saida: {PASTA_ENTREGA.resolve()}\n")

    if modo == "local":
        pptx, xlsx = gerar_apresentacao(modo_link="local")
        print(f"\n>>> Abra: {pptx.resolve()}")
        print(f">>> Planilha: {xlsx}")
        print("Versao e-mail: python -m scripts.gerar_apresentacao_diretoria email")
        try:
            import os
            os.startfile(str(pptx.resolve()))
            os.startfile(str(PASTA_ENTREGA.resolve()))
        except OSError as err:
            print(f"(Nao foi possivel abrir automaticamente: {err})")
    else:
        pptx, xlsx = gerar_apresentacao(modo_link="email")
        print(f"\nApresentação (e-mail): {pptx}")
        print(f"Planilha (anexar):     {xlsx}")
        print("Anexe os dois; quem receber salva na mesma pasta antes de clicar no link.")
